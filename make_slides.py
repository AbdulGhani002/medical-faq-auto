"""Generate MedFAQ_Slides.pptx from this script (no Markdown parsing).

  python make_slides.py

Produces ``MedFAQ_Slides.pptx`` in the repo root with 9 slides covering
scope, pipeline, dataset, models, evaluation, strengths, limitations,
future direction, and demo. Plain look, presentable directly.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HERE = Path(__file__).parent
OUT = HERE / "MedFAQ_Slides.pptx"

# 16:9 dimensions in EMUs (1 in = 914400 EMU). python-pptx defaults to
# 10 x 7.5; we keep that but render content within ~9.5 x 7.0 safe area.
INK = RGBColor(0x1C, 0x17, 0x17)
MUTED = RGBColor(0x78, 0x71, 0x6C)
ACCENT = RGBColor(0x05, 0x96, 0x69)  # emerald-600
BG_LIGHT = RGBColor(0xFA, 0xFA, 0xF9)


def _set_title(slide, text: str, size: int = 32) -> None:
    title = slide.shapes.title
    title.text = text
    run = title.text_frame.paragraphs[0].runs[0]
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = INK


def _add_text(slide, left, top, width, height, text: str,
              size: int = 14, bold: bool = False,
              color: RGBColor = INK, align=PP_ALIGN.LEFT) -> None:
    tx = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = ln
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def _add_bullets(slide, left, top, width, height,
                 bullets: list[tuple[str, str | None]],
                 size: int = 13) -> None:
    """Each bullet is (heading, body). body=None means single-line."""
    tx = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tx.text_frame
    tf.word_wrap = True
    for i, (head, body) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # Heading run (bold).
        r = p.add_run()
        r.text = "•  " + head
        r.font.size = Pt(size)
        r.font.bold = True
        r.font.color.rgb = INK
        if body:
            r2 = p.add_run()
            r2.text = "  — " + body
            r2.font.size = Pt(size - 1)
            r2.font.color.rgb = MUTED


def _add_table(slide, left, top, width, height,
               rows: list[list[str]], header: bool = True) -> None:
    n_rows, n_cols = len(rows), len(rows[0])
    tbl = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(left), Inches(top), Inches(width), Inches(height)
    ).table
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = val
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(12)
                    run.font.color.rgb = INK
                    if header and r == 0:
                        run.font.bold = True


def _accent_strip(slide) -> None:
    """Thin emerald strip at the top for visual consistency."""
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Emu(91440))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


# -------------------- build the deck --------------------


prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
TITLE = prs.slide_layouts[5]  # title only


def new_slide():
    s = prs.slides.add_slide(BLANK)
    _accent_strip(s)
    return s


# -------- Slide 1: title --------
s = new_slide()
_add_text(s, 0.7, 1.6, 12.0, 1.2,
          "MedFAQ",
          size=72, bold=True, color=INK)
_add_text(s, 0.7, 2.9, 12.0, 0.8,
          "Auto-Generated Medical FAQ Assistant",
          size=28, color=INK)
_add_text(s, 0.7, 3.8, 12.0, 0.6,
          "Retrieval-only chatbot · no LLM · every model built from scratch",
          size=18, color=MUTED)
_add_text(s, 0.7, 5.4, 12.0, 0.45,
          "Abdul Ghani · Anas Bhatti · Muhammad Salman · 2 teammates",
          size=14, color=INK)
_add_text(s, 0.7, 5.95, 12.0, 0.45,
          "NUTECH · 6th semester · github.com/AbdulGhani002/medical-faq-auto",
          size=14, color=MUTED)


# -------- Slide 2: Scope and requirements --------
s = new_slide()
_add_text(s, 0.5, 0.4, 12, 0.7,
          "Scope and requirements",
          size=32, bold=True, color=INK)
_add_text(s, 0.5, 1.25, 12, 0.5,
          "Problem · Patients ask the same questions every week.",
          size=14, color=MUTED)
_add_bullets(s, 0.7, 1.95, 12, 5,
             [
              ("Goal", "Retrieval chatbot that answers from an approved, clinician-reviewed FAQ index across three specialties."),
              ("Specialties", "Radiology · Physiotherapy · Cardiovascular."),
              ("Honesty", "Every reply is a real, human-written FAQ entry. No hallucination."),
              ("Hard constraint", "No LLM. No pretrained transformer. No sentence-transformers downloaded."),
              ("Hard constraint", "Every ML algorithm written by us (numpy or PyTorch tensors). ~3000 LoC."),
              ("Deploy", "Single `docker compose up` brings web + api + mongo live."),
              ("Pakistan context", "Roman-Urdu phonetic input · 1122 emergency banner · PHI scrub."),
             ], size=15)


# -------- Slide 3: Pipeline --------
s = new_slide()
_add_text(s, 0.5, 0.4, 12, 0.7,
          "Runtime pipeline",
          size=32, bold=True, color=INK)
_add_text(s, 0.5, 1.25, 12, 0.5,
          "Each user turn flows top-to-bottom in ~22 ms p50.",
          size=14, color=MUTED)
_add_bullets(s, 0.7, 1.95, 12, 5,
             [
              ("Preprocess", "PHI scrub → Roman-Urdu expand → coreference → context augment → spell-correct."),
              ("Understanding", "Tokeniser · POS (lexicon + HMM) · NER (dict + perceptron + CRF) · negation · sentiment · triage · intent (NB / MLP / Transformer)."),
              ("Retrieve (5 channels)", "TF-IDF word 0.20 · TF-IDF char 0.45 · BM25 0.25 · LSA 0.05 · PPMI 0.05."),
              ("Rerank", "PRF (Rocchio) → LTR (LogReg, 14 features) → MMR diversify."),
              ("Dialog", "Empathic opener · triage banner · clarification · disambiguation card · stem-overlap highlight."),
              ("Reply", "answer + matched FAQ + score breakdown + NLP analysis + dialog meta."),
              ("Offline pipeline", "Chat-log → segment → normalise → embed (TF-IDF) → cluster (K-Means) → select → polish → publish."),
             ], size=14)


# -------- Slide 4: Dataset --------
s = new_slide()
_add_text(s, 0.5, 0.4, 12, 0.7,
          "Dataset",
          size=32, bold=True, color=INK)
_add_text(s, 0.5, 1.25, 12, 0.5,
          "All data lives in `data/*.jsonl`. Regenerated by `data/build_jsonl.py`.",
          size=14, color=MUTED)
_add_table(s, 0.5, 1.85, 12.3, 4.7,
           [
            ["File", "Records", "Description"],
            ["faqs.jsonl", "169", "Hand-written clinician-style FAQs (3 specialties)"],
            ["intents.jsonl", "308", "Curated intent training, 13 classes incl. meta"],
            ["intents_augmented.jsonl", "1206", "6-way augmentation: synonyms / typos / prefixes"],
            ["eval_queries.jsonl", "116", "Held-out queries with expected keywords"],
            ["ner_examples.jsonl", "32", "BIO-tagged for CRF training"],
            ["sentiment_lexicon.jsonl", "96", "Polarity weights + intensifiers"],
            ["chat_logs.jsonl", "18", "Synthetic patient sessions"],
            ["kg_triples.jsonl", "518", "Auto-mined (subject, predicate, object) triples"],
           ])
_add_text(s, 0.5, 6.9, 12, 0.4,
          "Counts persist live in MongoDB on every chat match → popular FAQs surface to the top.",
          size=12, color=MUTED)


# -------- Slide 5: Models --------
s = new_slide()
_add_text(s, 0.5, 0.4, 12, 0.7,
          "Models — every one trained from scratch",
          size=30, bold=True, color=INK)
_add_text(s, 0.5, 1.2, 5.9, 0.4,
          "Classical numpy (myml/)",
          size=15, bold=True, color=INK)
_add_bullets(s, 0.55, 1.6, 6, 5,
             [
              ("TF-IDF", "word + char n-grams, sublinear TF"),
              ("BM25 Okapi", "Lucene-style smoothed IDF"),
              ("Naive Bayes", "Laplace smoothing"),
              ("Logistic Reg.", "batch GD + L2 + class weights"),
              ("Randomised SVD", "Halko-Martinsson-Tropp 2011"),
              ("K-Means", "k-means++ init"),
              ("MLP + backprop", "ReLU/softmax, momentum SGD"),
              ("CRF", "forward-backward + Viterbi"),
              ("Word2Vec", "skip-gram + neg sampling"),
              ("Dual encoder", "InfoNCE bi-encoder"),
              ("BPE tokenizer", "Sennrich 2016"),
             ], size=11)
_add_text(s, 6.8, 1.2, 6.0, 0.4,
          "Trainable transformer (PyTorch, GPU)",
          size=15, bold=True, color=INK)
_add_bullets(s, 6.85, 1.6, 6, 5,
             [
              ("Backend", "PyTorch tensors + autograd only"),
              ("No pretrained weights", "every layer written by us"),
              ("Arch.", "2 encoder blocks · 4 heads · d=64 · d_ff=128"),
              ("Embedding", "vocab-id → 64-d + sinusoidal pos enc"),
              ("Attention", "scaled dot-product · multi-head · pre-LN"),
              ("Pool / head", "masked mean-pool → softmax(13 classes)"),
              ("Training data", "1206 augmented intent examples"),
              ("Optimiser", "AdamW + cosine LR + grad clip"),
              ("Hardware", "Auto-detects CUDA; CPU fallback"),
              ("Result", "train acc 1.000, loss 0.0003"),
              ("Endpoint", "POST /nlp/torch_intent"),
             ], size=11)


# -------- Slide 6: Evaluation --------
s = new_slide()
_add_text(s, 0.5, 0.4, 12, 0.7,
          "Evaluation",
          size=32, bold=True, color=INK)
_add_text(s, 0.5, 1.25, 12, 0.5,
          "Retrieval evaluated on 116 held-out queries; transformer on training accuracy.",
          size=14, color=MUTED)
_add_table(s, 0.5, 1.95, 4.0, 2.3,
           [
            ["Metric", "Value"],
            ["Precision@1", "0.931"],
            ["MRR", "0.938"],
            ["Recall@5", "0.948"],
            ["Median latency", "22 ms"],
            ["p95 latency", "47 ms"],
           ])
_add_table(s, 4.9, 1.95, 8.0, 2.3,
           [
            ["Specialty", "n", "P@1", "Recall@5", "MRR"],
            ["Radiology", "36", "0.861", "0.889", "0.870"],
            ["Physiotherapy", "35", "0.971", "0.971", "0.971"],
            ["Cardiovascular", "45", "0.956", "0.978", "0.967"],
           ])
_add_text(s, 0.5, 4.6, 12, 0.5,
          "Transformer intent classifier:  train accuracy 1.000 · final CE loss 0.0003 over 40 epochs",
          size=14, color=INK)
_add_text(s, 0.5, 5.15, 12, 0.5,
          "Integration suite (api/final_test.py):  62 / 62 PASS  — covers chat, faq, admin, stats, nlp, frontend routes.",
          size=14, color=INK)
_add_text(s, 0.5, 5.85, 12, 0.5,
          "Live FAQ counters verified end-to-end:  counts start at 0, increment per match, persist in Mongo.",
          size=14, color=INK)


# -------- Slide 7: Strengths --------
s = new_slide()
_add_text(s, 0.5, 0.4, 12, 0.7,
          "Strengths",
          size=32, bold=True, color=INK)
_add_bullets(s, 0.6, 1.4, 12.4, 5.7,
             [
              ("No LLM dependency", "reproducible, auditable, no API costs, no hallucinations."),
              ("Every algorithm is ours", "drop-in replacements for sklearn, rank_bm25 and sentence-transformers."),
              ("Hybrid retrieval works", "P@1 0.93 across 116 held-out queries at 22 ms p50."),
              ("Multi-turn dialog", "coreference, slot filling, topic tracking, triage banner, empathic openers, disambiguation."),
              ("Pakistan-friendly", "Roman-Urdu phonetic mapping, 1122 emergency wording, PHI scrub at ingest."),
              ("One-command deploy", "docker compose up auto-seeds Mongo, builds API + Next.js."),
              ("Real-time analytics", "FAQ count increments live, persisted in Mongo, surfaced in /admin."),
              ("Trainable transformer", "GPU-ready PyTorch model trained from scratch on our own corpus."),
             ], size=15)


# -------- Slide 8: Limitations --------
s = new_slide()
_add_text(s, 0.5, 0.4, 12, 0.7,
          "Limitations",
          size=32, bold=True, color=INK)
_add_bullets(s, 0.6, 1.4, 12.4, 5.7,
             [
              ("Small corpus", "169 FAQs cover the most common questions but out-of-distribution queries fall through to best-guess."),
              ("Eval is approximate", "expected-keyword match in the answer can mark a correct FAQ as wrong when phrasing diverges."),
              ("Transformer overfits", "trained on 1206 augmented examples; held-out test of the classifier itself is still small."),
              ("NER is dictionary-first", "the CRF + structured perceptron variants improve coverage but training data is only 32 hand-annotated sentences."),
              ("Roman-Urdu is heuristic", "~80-word dictionary catches common cases; full transliteration would need a real parallel corpus."),
              ("Voice input is browser-native", "Web Speech API works in Chrome / Edge only."),
              ("Counts are at-most-once per turn", "no fraud / abuse heuristic — every successful match increments."),
             ], size=15)


# -------- Slide 9: Future direction --------
s = new_slide()
_add_text(s, 0.5, 0.4, 12, 0.7,
          "Future direction",
          size=32, bold=True, color=INK)
_add_bullets(s, 0.6, 1.4, 12.4, 5.7,
             [
              ("Grow the corpus", "500+ FAQs across more specialties with clinician validation."),
              ("Bi-encoder upgrade", "train the dual encoder on real clinician-rated query–FAQ pairs, not just keyword-derived positives."),
              ("Active learning loop", "low-confidence turns get queued for clinician review; approved Q/A pairs feed back into the corpus."),
              ("Real Urdu", "deeper transliteration trained on a parallel corpus; voice input in Urdu."),
              ("Clinical eval", "NDCG with clinician-rated relevance instead of a keyword heuristic."),
              ("GPU service", "dedicated container with CUDA wheel for heavier neural extensions."),
              ("Audit + governance", "version-controlled FAQ approval workflow, rollback, change-log."),
             ], size=15)


# -------- Slide 10: Demo --------
s = new_slide()
_add_text(s, 0.5, 0.4, 12, 0.7,
          "Live demo · how to run",
          size=32, bold=True, color=INK)
_add_text(s, 0.5, 1.25, 12, 0.5,
          "Single command brings the entire stack up.",
          size=14, color=MUTED)
_add_text(s, 0.6, 1.85, 12, 0.5,
          "docker compose up --build",
          size=18, bold=True, color=ACCENT)
_add_bullets(s, 0.6, 2.7, 12.4, 4.5,
             [
              ("http://localhost:3000/", "landing page · pick a specialty"),
              ("/chat/radiology", "full chatbot · NLP debug panel · voice input · reset / topic chip"),
              ("/playground", "type any sentence, see every NLP layer live"),
              ("/architecture", "system architecture diagram + technique table + metrics"),
              ("/admin", "live FAQ counts + chart-rich admin dashboard"),
              ("python api/train_transformer.py --verbose", "train the from-scratch transformer (GPU auto-detected)"),
              ("python api/eval.py", "regenerate EVALUATION_REPORT.md from 116 held-out queries"),
              ("python api/final_test.py", "62 / 62 integration suite"),
             ], size=14)


# ------------------ save ------------------

prs.save(OUT)
print(f"wrote {OUT.relative_to(HERE).as_posix()}  ({len(prs.slides)} slides)")
